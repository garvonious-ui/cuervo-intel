"""
Autostrat PDF-to-JSON Parser
=============================
Reads PDF exports (from Google Slides) of autostrat.ai reports,
extracts text via pdftotext, and outputs structured JSON matching
the templates in data/autostrat/{report_type}/.

Usage:
    from autostrat_parser import parse_all_pdfs
    results = parse_all_pdfs()  # scans data/autostrat/pdfs/
"""

from __future__ import annotations

import json
import os
import re
import subprocess

BASE_DIR = os.path.dirname(__file__)
AUTOSTRAT_DIR = os.path.join(BASE_DIR, "data", "autostrat")
PDF_DIR = os.path.join(AUTOSTRAT_DIR, "pdfs")

# Map title line to report type directory
REPORT_TYPE_MAP = {
    "tiktok hashtag analysis presentation": "tiktok_hashtags",
    "tiktok hashtag search analysis presentation": "tiktok_hashtags",
    "tiktok profile analysis presentation": "tiktok_profiles",
    "instagram profile analysis presentation": "instagram_profiles",
    "instagram profile presentation": "instagram_profiles",
    "instagram hashtag analysis presentation": "instagram_hashtags",
    "instagram hashtag search analysis presentation": "instagram_hashtags",
    "tiktok keyword analysis presentation": "tiktok_keywords",
    "tiktok keywords analysis presentation": "tiktok_keywords",
    "tiktok keyword search analysis presentation": "tiktok_keywords",
    "tiktok keywords search analysis presentation": "tiktok_keywords",
    "instagram keyword analysis presentation": "instagram_keywords",
    "instagram keyword search analysis presentation": "instagram_keywords",
    "instagram keywords analysis presentation": "instagram_keywords",
    "instagram keywords search analysis presentation": "instagram_keywords",
    "google news analysis presentation": "google_news",
    "google news presentation": "google_news",
}

# Section headings used as text delimiters (order matters for splitting)
SECTION_HEADINGS = [
    "Executive Summary",
    "Audience Profile",
    "Snapshot",
    "Creator Summary",
    "Hashtag Analysis",
    "Interesting Conversations",
    "Conversation Map",
    "Content Trends",
    "Brand Mentions",
    "In-Market Campaigns",
    "How to Win With This Audience",
    "Creator Archetypes",
    "Sponsorship Analysis",
    "Future Sponsorship Suggestions",
    "Engagement Analysis",
    "Posting Analysis",
    "Summary Statistics - By Post Type",
    "Summary Statistics - All Posts",
    "Summary Statistics",
    "Most / Least Liked",
    "Most / Least Comments",
    "Most / Least Engaged",
    "News Analysis",
    "News Trends",
    "News Topics",
    "Top Stories",
    "Competitor Coverage",
    "Trending Narratives",
    "SWOT Analysis",
    "Consideration Spaces",
    "Potential Actions",
    "Quotes",
    "Strategic Implications",
    "Appendix",
]


def extract_text_from_pdf(pdf_path: str) -> str:
    """Extract text from a PDF using pdftotext."""
    result = subprocess.run(
        ["pdftotext", pdf_path, "-"],
        capture_output=True, text=True, timeout=30
    )
    if result.returncode != 0:
        raise RuntimeError(f"pdftotext failed on {pdf_path}: {result.stderr}")
    return result.stdout


def extract_text_from_pptx(pptx_path: str) -> str:
    """Extract text from a PPTX file, producing output compatible with the PDF parser.

    Handles shape ordering (by position), group shapes (section headings),
    tables (NOPD, statistics, top posts, sponsorship), and injects a
    report-type header line so detect_report_type() works.
    """
    from pptx import Presentation as PptxPresentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    NOPD_LABELS = {"NEEDS", "OBJECTIONS", "DESIRES", "PAIN POINTS"}

    def _extract_shape_text(shape) -> list[str]:
        """Extract all text lines from a single shape."""
        lines = []
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                t = para.text.strip()
                if t:
                    lines.append(t)
        return lines

    def _extract_table_text(shape) -> list[str]:
        """Extract text from a table shape with smart formatting."""
        lines = []
        table = shape.table
        n_cols = len(table.columns)
        n_rows = len(list(table.rows))

        # Detect NOPD table: has rows where first cell is a NOPD label
        is_nopd = any(
            row.cells[0].text.strip() in NOPD_LABELS
            for row in table.rows
        )

        if is_nopd:
            # Output: LABEL\n\nitem1\n\nitem2\n\nitem3\n\n (blank-line separated)
            for row in table.rows:
                cells = [c.text.strip() for c in row.cells]
                first = cells[0]
                if first in NOPD_LABELS:
                    lines.append(first)
                else:
                    # Each non-empty cell is an individual item
                    for cell in cells:
                        if cell:
                            lines.append(cell)
                            lines.append("")  # blank line separator
        elif n_cols == 2:
            # Key-value table (stats, top posts): label then value
            for row in table.rows:
                cells = [c.text.strip() for c in row.cells]
                lines.append(cells[0])
                lines.append(cells[1])
        else:
            # Multi-column table (sponsorship): each cell on its own line
            for row in table.rows:
                for cell in row.cells:
                    t = cell.text.strip()
                    if t:
                        lines.append(t)

        return lines

    # Row-clustering threshold: shapes within this vertical distance
    # are considered the same row (50000 EMU ≈ 0.55 inch).
    ROW_SNAP = 50000

    prs = PptxPresentation(pptx_path)
    slide_blocks: list[str] = []

    for slide in prs.slides:
        # Collect all shapes with their position for sorting
        heading_lines: list[str] = []  # from group shapes (section headings)
        positioned: list[tuple[int, int, list[str]]] = []  # (top, left, lines)

        for shape in slide.shapes:
            top = shape.top or 0
            left = shape.left or 0

            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                # Group shapes contain section headings — collect separately
                for sub in shape.shapes:
                    heading_lines.extend(_extract_shape_text(sub))
            elif shape.has_table:
                tlines = _extract_table_text(shape)
                if tlines:
                    positioned.append((top, left, tlines))
            elif shape.has_text_frame:
                tlines = _extract_shape_text(shape)
                if tlines:
                    positioned.append((top, left, tlines))

        # Cluster shapes into rows by similar top positions, then sort by left
        positioned.sort(key=lambda x: (x[0], x[1]))
        if positioned:
            rows: list[list[tuple[int, int, list[str]]]] = []
            current_row: list[tuple[int, int, list[str]]] = [positioned[0]]
            for item in positioned[1:]:
                if abs(item[0] - current_row[0][0]) <= ROW_SNAP:
                    current_row.append(item)
                else:
                    rows.append(current_row)
                    current_row = [item]
            rows.append(current_row)
            # Sort within each row by left position
            positioned = []
            for row in rows:
                row.sort(key=lambda x: x[1])
                positioned.extend(row)

        # Assemble slide text with blank lines between shapes.
        # For the first slide, positioned content (identifier, date) goes before
        # headings (report type title). For all other slides, headings (section
        # names) go first.
        slide_lines: list[str] = []
        is_first_slide = (len(slide_blocks) == 0)

        if is_first_slide:
            # Put identifier + date before report type heading
            for _, _, lines in positioned:
                if slide_lines and slide_lines[-1] != "":
                    slide_lines.append("")
                slide_lines.extend(lines)
            if heading_lines:
                if slide_lines and slide_lines[-1] != "":
                    slide_lines.append("")
                slide_lines.extend(heading_lines)
        else:
            slide_lines.extend(heading_lines)
            for _, _, lines in positioned:
                if slide_lines and slide_lines[-1] != "":
                    slide_lines.append("")
                slide_lines.extend(lines)

        slide_blocks.append("\n".join(slide_lines))

    full_text = "\n\n".join(slide_blocks)

    # Detect report type from content and inject it into the header
    lower = full_text.lower()
    if "this instagram profile" in lower or "instagram profile" in lower:
        report_type_line = "Instagram Profile Analysis Presentation"
    elif "this tiktok profile" in lower or "tiktok profile" in lower:
        report_type_line = "TikTok Profile Analysis Presentation"
    elif "instagram hashtag" in lower:
        report_type_line = "Instagram Hashtag Analysis Presentation"
    elif "tiktok hashtag" in lower:
        report_type_line = "TikTok Hashtag Analysis Presentation"
    elif "tiktok keyword" in lower:
        report_type_line = "TikTok Keyword Analysis Presentation"
    elif "google news" in lower:
        report_type_line = "Google News Analysis Presentation"
    else:
        fname = os.path.basename(pptx_path).lower()
        if "profile" in fname:
            report_type_line = "Instagram Profile Analysis Presentation"
        elif "hashtag" in fname:
            report_type_line = "Instagram Hashtag Analysis Presentation"
        else:
            report_type_line = "Instagram Profile Analysis Presentation"

    # Only inject report type line if not already present in the first few lines
    first_lines = full_text[:500].lower()
    already_has_type = any(
        rt_key in first_lines for rt_key in REPORT_TYPE_MAP
    )

    if not already_has_type:
        # Insert report type line after the first line (identifier)
        first_nl = full_text.find("\n")
        if first_nl > 0:
            full_text = full_text[:first_nl] + "\n" + report_type_line + full_text[first_nl:]
        else:
            full_text = report_type_line + "\n" + full_text

    return full_text


def detect_report_type(text: str) -> tuple[str, str, str]:
    """Detect report type, identifier, and date from the first lines of text.

    Returns (report_type_dir, identifier, date_str).
    """
    lines = [ln.strip() for ln in text.split("\n") if ln.strip()]
    if len(lines) < 3:
        raise ValueError("PDF text too short to detect report type")

    identifier = lines[0].strip("'").strip()
    # Handle possessive forms like "Brand's TikTok Profile Analysis"
    if "'s " in identifier:
        identifier = identifier.split("'s ")[0]

    # Find the report type line
    report_type_dir = None
    for line in lines[:5]:
        key = line.strip().lower()
        if key in REPORT_TYPE_MAP:
            report_type_dir = REPORT_TYPE_MAP[key]
            break

    if not report_type_dir:
        # Try partial matching
        for line in lines[:5]:
            key = line.strip().lower()
            for pattern, rtype in REPORT_TYPE_MAP.items():
                if pattern in key or key in pattern:
                    report_type_dir = rtype
                    break
            if report_type_dir:
                break

    if not report_type_dir:
        raise ValueError(f"Could not detect report type from first lines: {lines[:5]}")

    # Find date (look for date-like patterns)
    date_str = ""
    for line in lines[:8]:
        # Match "January 06, 2026" or "2025-07-04" or "December 10, 2025"
        m = re.search(r'(\d{4}-\d{2}-\d{2})', line)
        if m:
            date_str = m.group(1)
            break
        m = re.search(r'(\w+ \d{1,2},?\s*\d{4})', line)
        if m:
            date_str = m.group(1)
            break

    return report_type_dir, identifier, date_str


def split_into_sections(text: str) -> dict[str, str]:
    """Split text into named sections using heading delimiters."""
    sections = {}
    # Remove the boilerplate "How to use this deck" section
    boilerplate_end = text.find("Autostrat Team")
    if boilerplate_end > 0:
        # Find the next line break after boilerplate
        next_nl = text.find("\n", boilerplate_end)
        if next_nl > 0:
            text = text[:text.find("How to use this deck")] + text[next_nl:]

    # Build a list of (position, heading_name) for all headings found
    found = []
    for heading in SECTION_HEADINGS:
        # Use word-boundary matching to avoid partial matches
        pattern = re.escape(heading)
        for m in re.finditer(pattern, text):
            found.append((m.start(), heading, m.end()))

    # Sort by position
    found.sort(key=lambda x: x[0])

    # Remove substring matches: when two headings overlap (one contains the other),
    # keep only the longer heading at that position.
    filtered = []
    for pos, name, end in found:
        # Check if this match is contained within a longer heading match
        is_substring = False
        for opos, oname, oend in found:
            if oname == name:
                continue
            # If another heading starts at same position or contains this one
            if opos <= pos and oend >= end and len(oname) > len(name):
                is_substring = True
                break
        if not is_substring:
            filtered.append((pos, name, end))
    found = filtered

    # Deduplicate — if same heading appears multiple times, keep first before Appendix
    # and any after Appendix separately
    appendix_pos = None
    for pos, name, end in found:
        if name == "Appendix":
            appendix_pos = pos
            break

    # Extract sections
    for i, (pos, name, end) in enumerate(found):
        # Skip duplicate headings that appear in the appendix
        if appendix_pos and pos > appendix_pos and name != "Appendix":
            # Store as appendix subsection
            next_pos = found[i + 1][0] if i + 1 < len(found) else len(text)
            key = f"appendix_{name}"
            sections[key] = text[end:next_pos].strip()
            continue

        next_pos = found[i + 1][0] if i + 1 < len(found) else len(text)
        section_text = text[end:next_pos].strip()
        sections[name] = section_text

    return sections


def parse_number(s: str) -> int | float:
    """Parse a number string, handling commas and percentages."""
    s = s.strip().replace(",", "").replace("%", "").strip()
    try:
        if "." in s:
            return float(s)
        return int(s)
    except (ValueError, TypeError):
        return 0


def parse_nopd(text: str) -> dict:
    """Parse Needs/Objections/Desires/Pain Points from audience profile text."""
    result = {"summary": "", "needs": [], "objections": [], "desires": [], "pain_points": []}

    # Split on the NOPD labels
    parts = re.split(r'\n(NEEDS|OBJECTIONS|DESIRES|PAIN POINTS)\n', text)

    if parts:
        result["summary"] = parts[0].strip()

    current_key = None
    key_map = {
        "NEEDS": "needs",
        "OBJECTIONS": "objections",
        "DESIRES": "desires",
        "PAIN POINTS": "pain_points",
    }

    for part in parts[1:]:
        part = part.strip()
        if part in key_map:
            current_key = key_map[part]
        elif current_key:
            # Split into individual items on double newlines or clear paragraph breaks
            items = re.split(r'\n\n+', part)
            for item in items:
                item = item.strip()
                if item and len(item) > 5:
                    result[current_key].append(item)

    return result


def parse_snapshot(text: str) -> dict:
    """Parse snapshot metrics (followers, following, avg likes, etc.).

    Handles both PDF layout (labels grouped then values grouped) and
    PPTX layout (labels and values may interleave in any order).
    Strategy: find each label, then take the nearest following numeric line.
    """
    snapshot = {
        "followers": 0, "following": 0,
        "avg_likes": 0, "avg_comments": 0, "avg_engagement_rate": 0
    }

    lines = [ln.strip() for ln in text.split("\n")]
    LABEL_MAP = {
        "followers": "followers",
        "following": "following",
        "avg likes": "avg_likes",
        "avg comments": "avg_comments",
        "avg engagement rate": "avg_engagement_rate",
    }

    # Find indices of each label
    label_positions: list[tuple[int, str]] = []
    for i, line in enumerate(lines):
        low = line.lower()
        if low in LABEL_MAP:
            label_positions.append((i, LABEL_MAP[low]))

    # For each label, find the nearest following numeric line
    # that isn't claimed by another label closer to it.
    # Extended range (12) handles PPTX text where blank lines separate shapes.
    claimed: set[int] = set()
    for pos, key in label_positions:
        for j in range(pos + 1, min(pos + 12, len(lines))):
            if j in claimed:
                continue
            line = lines[j].strip()
            if not line:
                continue
            # Skip other label lines
            if line.lower() in LABEL_MAP:
                continue
            if re.match(r'^[\d,]+\.?\d*\s*%?$', line):
                snapshot[key] = parse_number(line)
                claimed.add(j)
                break

    return snapshot


def parse_creator_summary(text: str) -> dict:
    """Parse creator summary section."""
    result = {
        "search_purpose": "",
        "topline": "",
        "what_it_means": "",
        "common_themes": [],
        "what_hits": "",
        "what_misses": "",
    }

    # Split by known sub-headings
    sub_headings = ["Search Purpose", "Topline", "What it Means for You",
                    "What it Means", "Common Themes and Topics",
                    "Common Themes", "What Hits", "What Misses"]

    parts = {}
    current = "preamble"
    current_text = []

    for line in text.split("\n"):
        stripped = line.strip()
        matched = False
        for heading in sub_headings:
            if stripped == heading:
                parts[current] = "\n".join(current_text).strip()
                current = heading
                current_text = []
                matched = True
                break
        if not matched:
            current_text.append(line)

    parts[current] = "\n".join(current_text).strip()

    result["search_purpose"] = parts.get("Search Purpose", "")
    result["topline"] = parts.get("Topline", parts.get("preamble", ""))
    result["what_it_means"] = parts.get("What it Means for You",
                                         parts.get("What it Means", ""))
    themes_text = parts.get("Common Themes and Topics",
                            parts.get("Common Themes", ""))
    if themes_text:
        result["common_themes"] = [t.strip() for t in re.split(r'\n\n+', themes_text)
                                   if t.strip() and len(t.strip()) > 5]
    result["what_hits"] = parts.get("What Hits", "")
    result["what_misses"] = parts.get("What Misses", "")

    return result


def parse_how_to_win(text: str) -> dict:
    """Parse How to Win section.

    Structure: summary, "Audience Verbatims" label, Territory 1-5 text,
    then actual verbatim quotes after the last territory.
    """
    result = {"summary": "", "territories": [], "audience_verbatims": []}

    # Split on Territory markers
    territory_pattern = r'Territory \d+'
    parts = re.split(territory_pattern, text)

    if parts:
        first = parts[0].strip()
        # Remove "Audience Verbatims" label from summary
        verbatim_marker = "Audience Verbatims"
        if verbatim_marker in first:
            idx = first.index(verbatim_marker)
            result["summary"] = first[:idx].strip()
        else:
            result["summary"] = first

    # Parse territories — each part between Territory N markers
    for part in parts[1:]:
        part = part.strip()
        if part and len(part) > 10:
            paras = re.split(r'\n\n+', part)
            if paras:
                result["territories"].append(paras[0].strip())

    # Verbatims come AFTER the last territory text
    # Find the last territory's end position, then grab everything after
    last_territory_match = None
    for m in re.finditer(territory_pattern, text):
        last_territory_match = m

    if last_territory_match:
        after_last = text[last_territory_match.end():].strip()
        # Skip the territory text (first paragraph) to get to verbatims
        paras = re.split(r'\n\n+', after_last)
        # First para is the territory text; remaining are verbatims
        for para in paras[1:]:
            para = para.strip()
            if para and len(para) > 8:
                # Split multi-line quotes into individual verbatims
                for line in para.split("\n"):
                    line = line.strip()
                    if line and len(line) > 8:
                        result["audience_verbatims"].append(line)

    return result


def parse_sponsorships(text: str) -> dict:
    """Parse sponsorship analysis section."""
    result = {
        "summary": "",
        "integration_summary": "",
        "categories": [],
        "companies": [],
    }

    sub_headings = ["Sponsorship Summary", "Current Categories",
                    "Integration Summary", "Current Companies"]
    parts = {}
    current = "preamble"
    current_text = []

    for line in text.split("\n"):
        stripped = line.strip()
        matched = False
        for heading in sub_headings:
            if stripped == heading:
                parts[current] = "\n".join(current_text).strip()
                current = heading
                current_text = []
                matched = True
                break
        if not matched:
            current_text.append(line)
    parts[current] = "\n".join(current_text).strip()

    result["summary"] = parts.get("Sponsorship Summary", parts.get("preamble", ""))
    result["integration_summary"] = parts.get("Integration Summary", "")

    cats = parts.get("Current Categories", "")
    if cats:
        result["categories"] = [c.strip() for c in cats.split("\n")
                                if c.strip() and 2 < len(c.strip()) < 40
                                and not c.strip().endswith(".")]

    comps = parts.get("Current Companies", "")
    if comps:
        result["companies"] = [c.strip() for c in comps.split("\n")
                               if c.strip() and 1 < len(c.strip()) < 40
                               and not c.strip().endswith(".")]

    return result


def parse_future_sponsorships(text: str) -> list[dict]:
    """Parse Future Sponsorship Suggestions from multi-column PDF layout.

    The section has a summary paragraph, then N category names, short
    description snippets, "Why it Works" × N, Why text (interleaved),
    "How to Activate" × N, and How text (may be interleaved or sequential).
    """
    lines_all = text.split("\n")
    lines = [ln.strip() for ln in lines_all if ln.strip()]

    # Count Why/How labels
    why_label = {"Why it Works", "Why it works"}
    how_label = {"How to Activate", "How to activate"}
    why_indices = [i for i, ln in enumerate(lines) if ln in why_label]
    how_indices = [i for i, ln in enumerate(lines) if ln in how_label]

    n_cats = max(len(why_indices), len(how_indices), 1)

    if n_cats <= 1:
        # Single category or no labels — fall back to simple parsing
        return _parse_future_sponsorships_simple(text)

    # ── Summary: long lines at the start ────────────────────────────────
    summary_parts = []
    summary_end = 0
    for i, ln in enumerate(lines):
        if len(ln) > 60:
            summary_parts.append(ln)
            summary_end = i + 1
        elif summary_parts:
            break

    # ── Category names: use blank-line-separated groups from raw text ────
    # This correctly handles wrapped names like "Celebrity-Driven Cultural\nCollaborations"
    # which appear as consecutive lines without a blank line between them.

    # Find summary end and first Why label in raw text (lines_all)
    summary_end_raw = 0
    in_summary = False
    for i, ln in enumerate(lines_all):
        if len(ln.strip()) > 60:
            in_summary = True
            summary_end_raw = i + 1
        elif in_summary and ln.strip():
            break

    first_why_raw = len(lines_all)
    for i, ln in enumerate(lines_all):
        if ln.strip() in why_label:
            first_why_raw = i
            break

    # Parse blank-line-separated groups between summary and first Why
    groups: list[str] = []
    current_group: list[str] = []
    for i in range(summary_end_raw, first_why_raw):
        ln = lines_all[i].strip()
        if ln:
            current_group.append(ln)
        elif current_group:
            groups.append(" ".join(current_group))
            current_group = []
    if current_group:
        groups.append(" ".join(current_group))

    # Separate names (short, no trailing period) from description snippets
    cat_names = []
    desc_lines = []
    for g in groups:
        if len(cat_names) < n_cats and not g.endswith((".", "!", "?")) and len(g) < 50:
            cat_names.append(g)
        else:
            desc_lines.append(g)

    # ── Description snippets: de-interleave by n_cats ───────────────────
    desc_texts = [""] * n_cats
    if desc_lines:
        cols = deinterleave_columns(desc_lines, n_cats)
        for ci, col_lines in enumerate(cols):
            desc_texts[ci] = " ".join(col_lines)

    # ── Why text: between last Why label and first How label ────────────
    last_why = why_indices[-1] if why_indices else -1
    first_how = how_indices[0] if how_indices else len(lines)

    why_text_lines = []
    for i in range(last_why + 1, first_how):
        ln = lines[i]
        if ln not in why_label:
            why_text_lines.append(ln)

    why_texts = [""] * n_cats
    if why_text_lines:
        cols = deinterleave_columns(why_text_lines, n_cats)
        for ci, col_lines in enumerate(cols):
            why_texts[ci] = " ".join(col_lines)

    # ── How to Activate text: after last How label ──────────────────────
    last_how = how_indices[-1] if how_indices else len(lines) - 1
    # Use raw text (preserving blank lines) for paragraph splitting
    # Find position of last How label in the original text
    raw_after_how = ""
    how_count = 0
    for i, ln in enumerate(lines_all):
        if ln.strip() in how_label:
            how_count += 1
            if how_count == len(how_indices):
                raw_after_how = "\n".join(lines_all[i + 1:])
                break

    how_blocks_raw = [b.strip() for b in re.split(r'\n\n+', raw_after_how)
                      if b.strip() and b.strip() not in how_label]
    # Filter out blocks from the next section (Engagement Analysis, etc.)
    how_blocks = []
    for block in how_blocks_raw:
        first_line = block.split("\n")[0].strip()
        if first_line in ("Engagement Analysis", "Summary", "Snapshot",
                          "Posting Analysis"):
            break
        how_blocks.append(block)

    how_items: list[list[str]] = [[] for _ in range(n_cats)]
    if how_blocks:
        if len(how_blocks) == n_cats:
            # Direct assignment: one block per category (cleanest case)
            for ci, block in enumerate(how_blocks):
                items = _split_action_items(block)
                how_items[ci] = items
        else:
            # Variable layout — join all and de-interleave by n_cats
            all_how_lines = []
            for block in how_blocks:
                for ln in block.split("\n"):
                    ln = ln.strip()
                    if ln and ln not in how_label:
                        all_how_lines.append(ln)
            if all_how_lines:
                cols = deinterleave_columns(all_how_lines, n_cats)
                for ci, col_lines in enumerate(cols):
                    how_items[ci] = _split_action_items(" ".join(col_lines))

    # ── Build suggestions ───────────────────────────────────────────────
    suggestions = []
    for ci in range(n_cats):
        name = cat_names[ci] if ci < len(cat_names) else f"Category {ci + 1}"
        why = why_texts[ci] if ci < len(why_texts) else ""
        if ci < len(desc_texts) and desc_texts[ci]:
            why = desc_texts[ci] + " " + why
        suggestions.append({
            "category": name,
            "why_it_works": why.strip(),
            "how_to_activate": how_items[ci] if ci < len(how_items) else [],
        })

    return suggestions


def _split_action_items(text: str) -> list[str]:
    """Split a text block into individual action items on sentence boundaries."""
    items = []
    for sentence in re.split(r'(?<=\.)\s+(?=[A-Z])', text):
        sentence = " ".join(sentence.split()).strip()
        if sentence and len(sentence) > 10:
            items.append(sentence)
    return items


def _parse_future_sponsorships_simple(text: str) -> list[dict]:
    """Fallback parser for single-category or unlabeled sponsorship sections."""
    suggestions = []
    current_cat = None
    why_text: list[str] = []
    how_text: list[str] = []
    in_why = False
    in_how = False

    for line in text.split("\n"):
        stripped = line.strip()
        if not stripped:
            continue
        if "Why it Works" in stripped or "Why it works" in stripped:
            if current_cat and (why_text or how_text):
                suggestions.append({
                    "category": current_cat,
                    "why_it_works": " ".join(why_text).strip(),
                    "how_to_activate": [h for h in how_text if h],
                })
                why_text, how_text = [], []
            in_why = True
            in_how = False
            continue
        if "How to Activate" in stripped or "How to activate" in stripped:
            in_why = False
            in_how = True
            continue

        if in_why:
            why_text.append(stripped)
        elif in_how:
            how_text.append(stripped)
        elif len(stripped) < 60 and not stripped.endswith("."):
            if current_cat and (why_text or how_text):
                suggestions.append({
                    "category": current_cat,
                    "why_it_works": " ".join(why_text).strip(),
                    "how_to_activate": [h for h in how_text if h],
                })
                why_text, how_text = [], []
            current_cat = stripped

    if current_cat and (why_text or how_text):
        suggestions.append({
            "category": current_cat,
            "why_it_works": " ".join(why_text).strip(),
            "how_to_activate": [h for h in how_text if h],
        })

    return suggestions


def parse_content_trends(text: str) -> list[dict]:
    """Parse content trends section."""
    trends = []
    # Trends come as title + description paragraphs
    paragraphs = re.split(r'\n\n+', text)
    i = 0
    while i < len(paragraphs):
        para = paragraphs[i].strip()
        if not para:
            i += 1
            continue
        # Short lines are likely titles, longer ones are descriptions
        if len(para) < 80 and i + 1 < len(paragraphs):
            trends.append({
                "trend": para,
                "description": paragraphs[i + 1].strip(),
            })
            i += 2
        else:
            # Might be a combined title+description
            lines = para.split("\n")
            if len(lines) >= 2:
                trends.append({
                    "trend": lines[0].strip(),
                    "description": " ".join(lines[1:]).strip(),
                })
            i += 1

    return trends


def parse_creator_archetypes(text: str) -> list[dict]:
    """Parse creator archetypes section."""
    archetypes = []

    # Look for "The ..." patterns as archetype names
    # Pattern: "The ASMR Mechanic", "The Relatable DIYer", etc.
    the_pattern = r'(The [A-Z][^\n]{3,40})'
    matches = list(re.finditer(the_pattern, text))

    if not matches:
        # Try splitting on "Appeal" markers
        parts = re.split(r'\nAppeal\n', text)
        for part in parts:
            part = part.strip()
            if part and len(part) > 20:
                lines = [ln.strip() for ln in part.split("\n") if ln.strip()]
                if lines:
                    archetypes.append({
                        "archetype": lines[0],
                        "description": " ".join(lines[1:3]) if len(lines) > 1 else "",
                        "appeal": "",
                        "examples": [],
                    })
        return archetypes

    for i, m in enumerate(matches):
        name = m.group(1).strip()
        start = m.end()
        end = matches[i + 1].start() if i + 1 < len(matches) else len(text)
        block = text[start:end].strip()

        arch = {"archetype": name, "description": "", "appeal": "", "examples": []}

        # Look for Appeal and Examples sub-sections
        if "Appeal" in block:
            parts = block.split("Appeal")
            arch["description"] = parts[0].strip()
            if len(parts) > 1:
                rest = parts[1]
                if "Examples" in rest:
                    appeal_part, examples_part = rest.split("Examples", 1)
                    arch["appeal"] = appeal_part.strip()
                    arch["examples"] = [ex.strip() for ex in examples_part.strip().split("\n")
                                        if ex.strip() and len(ex.strip()) > 5]
                else:
                    arch["appeal"] = rest.strip()
        else:
            arch["description"] = block

        archetypes.append(arch)

    return archetypes


def parse_brand_mentions(text: str) -> list[dict]:
    """Parse brand mentions section."""
    mentions = []

    # Look for brand name headers followed by context/sentiment/verbatims
    # Split on patterns like "Context\n" and "Sentiment\n"
    sub_headings = ["Context", "Reception", "Sentiment", "Verbatims"]

    # First try to find brand name blocks
    current_brand = None
    current_data = {"brand": "", "context": "", "sentiment": "", "reception": "", "verbatims": []}
    current_field = None

    for line in text.split("\n"):
        stripped = line.strip()
        if not stripped:
            continue

        if stripped in sub_headings:
            current_field = stripped.lower()
            continue

        # Check if this is a brand name (short, capitalized, not a sentence)
        if (len(stripped) < 30 and not stripped.endswith(".")
                and stripped[0].isupper() and current_field is None):
            if current_brand and current_data.get("context"):
                mentions.append(dict(current_data))
            current_brand = stripped
            current_data = {"brand": stripped, "context": "", "sentiment": "",
                            "reception": "", "verbatims": []}
            current_field = None
            continue

        if current_field == "context":
            current_data["context"] += " " + stripped
        elif current_field == "sentiment":
            current_data["sentiment"] += " " + stripped
        elif current_field == "reception":
            current_data["reception"] += " " + stripped
        elif current_field == "verbatims":
            current_data["verbatims"].append(stripped)
        elif current_brand and not current_field:
            # Probably part of the brand description
            current_data["context"] += " " + stripped

    if current_brand and current_data.get("context"):
        mentions.append(dict(current_data))

    # Clean up
    for m in mentions:
        m["context"] = m["context"].strip()
        m["sentiment"] = m["sentiment"].strip()
        m["reception"] = m["reception"].strip()

    return mentions


def parse_hashtag_analysis(text: str) -> dict:
    """Parse hashtag analysis section with findings/opportunities/gaps/actions.

    PDF 2-column layout produces paired headings on same row:
      Key Findings  |  Opportunities     → items interleaved [finding, opp, ...]
      Gaps/Risks    |  Strategic Actions  → items interleaved [gap, action, ...]
    We detect this pattern and de-interleave.
    """
    result = {
        "summary": "",
        "related_hashtags": [],
        "key_findings": [],
        "opportunities": [],
        "gaps_risks_unmet_needs": [],
        "strategic_actions": [],
    }

    # Locate headings by position
    heading_positions = {}
    for heading in ["Key Findings", "Opportunities",
                    "Gaps, Risks or Unmet Needs", "Strategic Actions"]:
        idx = text.find(heading)
        if idx >= 0:
            heading_positions[heading] = idx

    if not heading_positions:
        result["summary"] = text.strip()
        return result

    # Sort headings by position
    sorted_headings = sorted(heading_positions.items(), key=lambda x: x[1])

    # Get summary (text before first heading)
    first_pos = sorted_headings[0][1]
    if first_pos > 0:
        result["summary"] = text[:first_pos].strip()

    # Determine heading groups — paired headings that appear close together
    # are from the same 2-column row in the PDF
    groups = []
    i = 0
    while i < len(sorted_headings):
        name, pos = sorted_headings[i]
        # Check if the next heading is close (on same row = within ~80 chars)
        if i + 1 < len(sorted_headings):
            next_name, next_pos = sorted_headings[i + 1]
            gap_text = text[pos + len(name):next_pos].strip()
            if len(gap_text) < 30:
                # These two headings are a 2-column pair
                groups.append((name, next_name))
                i += 2
                continue
        groups.append((name,))
        i += 1

    # Extract content for each group
    for gi, group in enumerate(groups):
        # Find start of content (after all headings in this group)
        last_heading = group[-1]
        content_start = heading_positions[last_heading] + len(last_heading)

        # Find end of content (start of next group, or end of major section)
        if gi + 1 < len(groups):
            next_group_first = groups[gi + 1][0]
            content_end = heading_positions[next_group_first]
        else:
            # End at next major section or text end
            content_end = len(text)
            for boundary in ["Interesting Conversations", "Content Trends",
                             "Brand Mentions", "How to Win"]:
                bpos = text.find(boundary, content_start)
                if content_start < bpos < content_end:
                    content_end = bpos

        section_text = text[content_start:content_end].strip()
        items = [item.strip() for item in re.split(r'\n\n+', section_text)
                 if item.strip() and len(item.strip()) > 10]

        if len(group) == 2:
            # De-interleave: even indices = left column, odd = right column
            left_key = {
                "Key Findings": "key_findings",
                "Gaps, Risks or Unmet Needs": "gaps_risks_unmet_needs",
            }.get(group[0], "key_findings")
            right_key = {
                "Opportunities": "opportunities",
                "Strategic Actions": "strategic_actions",
            }.get(group[1], "opportunities")

            for idx, item in enumerate(items):
                if idx % 2 == 0:
                    result[left_key].append(item)
                else:
                    result[right_key].append(item)
        else:
            # Single heading — all items go to that key
            key_map = {
                "Key Findings": "key_findings",
                "Opportunities": "opportunities",
                "Gaps, Risks or Unmet Needs": "gaps_risks_unmet_needs",
                "Strategic Actions": "strategic_actions",
            }
            key = key_map.get(group[0], "key_findings")
            result[key] = items

    return result


def parse_interesting_conversations(text: str) -> list[dict]:
    """Parse interesting conversations section."""
    conversations = []

    # Look for "Conversation N" patterns
    conv_pattern = r'Conversation \d+'
    parts = re.split(conv_pattern, text)

    # First part is the summary
    summary = parts[0].strip() if parts else ""

    for part in parts[1:]:
        part = part.strip()
        if not part:
            continue
        # First line is title, rest is description
        lines = [ln.strip() for ln in part.split("\n") if ln.strip()]
        if lines:
            conversations.append({
                "title": lines[0],
                "description": " ".join(lines[1:]) if len(lines) > 1 else "",
            })

    return conversations


def parse_conversation_map(text: str) -> dict:
    """Parse conversation map section."""
    result = {
        "summary": "",
        "relationship_analysis": "",
        "overarching_patterns": [],
        "action_opportunities": [],
    }

    sub_headings = {
        "Conversation Map Analysis": None,
        "Relationship Analysis": "relationship_analysis",
        "Overarching Patterns": "overarching_patterns",
        "Conversation Action Opportunities": "action_opportunities",
    }

    current_key = "summary"
    current_text = []

    for line in text.split("\n"):
        stripped = line.strip()
        matched = False
        for heading, key in sub_headings.items():
            if stripped == heading:
                # Save current
                if current_key == "summary":
                    result["summary"] = "\n".join(current_text).strip()
                elif current_key in ("relationship_analysis",):
                    result[current_key] = "\n".join(current_text).strip()
                elif current_key in ("overarching_patterns", "action_opportunities"):
                    items = [t.strip() for t in "\n".join(current_text).strip().split("\n\n")
                             if t.strip() and len(t.strip()) > 10]
                    result[current_key] = items

                current_key = key or "summary"
                current_text = []
                matched = True
                break
        if not matched:
            current_text.append(line)

    # Save last section
    if current_key and current_text:
        text_joined = "\n".join(current_text).strip()
        if current_key in ("overarching_patterns", "action_opportunities"):
            result[current_key] = [t.strip() for t in text_joined.split("\n\n")
                                   if t.strip() and len(t.strip()) > 10]
        elif current_key in result and isinstance(result.get(current_key), str):
            result[current_key] = text_joined

    return result


def parse_in_market_campaigns(text: str) -> list[dict]:
    """Parse in-market campaigns section."""
    campaigns = []
    # Similar to content trends — title + description pairs
    paragraphs = re.split(r'\n\n+', text)
    i = 0
    while i < len(paragraphs):
        para = paragraphs[i].strip()
        if not para:
            i += 1
            continue
        if len(para) < 80 and i + 1 < len(paragraphs):
            campaigns.append({
                "campaign": para,
                "description": paragraphs[i + 1].strip(),
            })
            i += 2
        else:
            lines = para.split("\n")
            if len(lines) >= 2:
                campaigns.append({
                    "campaign": lines[0].strip(),
                    "description": " ".join(lines[1:]).strip(),
                })
            i += 1
    return campaigns


def parse_news_trends(text: str) -> list[dict]:
    """Parse News Trends section into trending narratives.

    PDF 2-column layout produces: [title1, title2, desc1, desc2, title3, title4, desc3, desc4].
    Titles are short (<120 chars), descriptions are longer.
    """
    paragraphs = [p.strip() for p in re.split(r'\n\n+', text) if p.strip()]

    titles = []
    descriptions = []

    for para in paragraphs:
        clean = " ".join(para.split())
        if len(clean) < 120:
            titles.append(clean)
        else:
            descriptions.append(clean)

    narratives = []
    for i in range(min(len(titles), len(descriptions))):
        narratives.append({
            "narrative": titles[i],
            "description": descriptions[i],
            "brands_involved": [],
        })

    # If there are extra titles without descriptions, still include them
    for i in range(len(descriptions), len(titles)):
        narratives.append({
            "narrative": titles[i],
            "description": "",
            "brands_involved": [],
        })

    return narratives


def parse_news_topics(text: str) -> dict:
    """Parse News Topics section into topic titles and detailed findings.

    Same 2-column layout: short topic titles, then longer finding descriptions.
    """
    paragraphs = [p.strip() for p in re.split(r'\n\n+', text) if p.strip()]

    topics = []
    findings = []

    for para in paragraphs:
        clean = " ".join(para.split())
        if len(clean) < 120:
            topics.append(clean)
        else:
            findings.append(clean)

    return {"topics": topics, "findings": findings}


def parse_swot_overview(text: str) -> dict:
    """Parse SWOT overview from Google News report.

    2-column layout: Strengths/Weaknesses header, then interleaved S/W items,
    then interleaved O/T items, then Opportunities/Threats label at bottom.
    Even-indexed items = left column, odd-indexed = right column.
    """
    result = {"strengths": [], "weaknesses": [], "opportunities": [], "threats": []}

    # Remove header labels
    clean = text
    clean = re.sub(r'\bStrengths\b', '', clean)
    clean = re.sub(r'\bWeaknesses\b', '', clean)

    # Find the Opportunities/Threats label line (appears at the bottom)
    lines = clean.split("\n")
    opp_label_idx = None
    for i, line in enumerate(lines):
        stripped = line.strip()
        if stripped == "Opportunities" or ("Opportunities" in stripped and "Threats" in stripped):
            opp_label_idx = i
            break

    if opp_label_idx is not None:
        # Remove the label line and anything after
        clean = "\n".join(lines[:opp_label_idx])

    clean = re.sub(r'\bOpportunities\b', '', clean)
    clean = re.sub(r'\bThreats\b', '', clean)

    # Get all items as paragraphs
    items = [p.strip() for p in re.split(r'\n\n+', clean) if p.strip() and len(p.strip()) > 20]

    # First half = S/W interleaved, second half = O/T interleaved
    half = len(items) // 2
    sw_items = items[:half]
    ot_items = items[half:]

    for i, item in enumerate(sw_items):
        c = " ".join(item.split())
        if i % 2 == 0:
            result["strengths"].append(c)
        else:
            result["weaknesses"].append(c)

    for i, item in enumerate(ot_items):
        c = " ".join(item.split())
        if i % 2 == 0:
            result["opportunities"].append(c)
        else:
            result["threats"].append(c)

    return result


def parse_potential_actions(text: str) -> list[str]:
    """Parse Potential Actions section into a list of action strings."""
    actions = []
    for para in re.split(r'\n\n+', text):
        para = para.strip()
        if not para or len(para) < 15:
            continue
        # Skip labels
        if para in ("Short Term", "Long Term", "Potential Actions"):
            continue
        actions.append(" ".join(para.split()))
    return actions


def parse_news_quotes(text: str) -> list[str]:
    """Parse Quotes section into individual quote strings.

    Quotes are separated by single newlines where each quote may wrap across
    multiple lines. Split on boundaries where a line ends with punctuation
    and the next line starts a new sentence (capital letter, quote mark, or number).
    """
    # First try double-newline split
    paras = [p.strip() for p in re.split(r'\n\n+', text) if p.strip()]

    quotes = []
    for para in paras:
        # Split multi-quote paragraphs on sentence boundaries
        lines = para.split("\n")
        current = []
        for line in lines:
            line = line.strip()
            if not line:
                continue
            if (current and
                    current[-1].rstrip().endswith((".","!","?",'"',"'")) and
                    len(line) > 5 and
                    (line[0].isupper() or line[0] == '"' or line[0] == "'")):
                # This looks like the start of a new quote
                joined = " ".join(current)
                if len(joined) > 20:
                    quotes.append(joined)
                current = [line]
            else:
                current.append(line)
        if current:
            joined = " ".join(current)
            if len(joined) > 20:
                quotes.append(joined)

    return quotes


def parse_news_statistics(text: str) -> list[str]:
    """Parse Statistics section into individual stat strings.

    Stats are separated by newlines but may wrap. A new stat typically starts
    with a capital letter or brand name and contains a number.
    """
    lines = [ln.strip() for ln in text.split("\n") if ln.strip()]
    stats = []
    current = []

    for line in lines:
        if not line:
            continue
        # A new stat starts with a capital letter and builds on a fresh sentence
        if (current and
                line[0].isupper() and
                current[-1].rstrip().endswith((".", "!", "?", "%"))):
            joined = " ".join(current)
            if len(joined) > 15:
                stats.append(joined)
            current = [line]
        else:
            current.append(line)

    if current:
        joined = " ".join(current)
        if len(joined) > 15:
            stats.append(joined)

    return stats


def parse_executive_summary(text: str) -> dict:
    """Parse executive summary section."""
    result = {"overview": "", "key_insights": [], "search_term": "", "search_purpose": ""}

    # Look for "What You Searched" and "Why You're Searching"
    if "What You Searched" in text:
        parts = text.split("What You Searched")
        if len(parts) > 1:
            rest = parts[1]
            if "Why You're Searching" in rest:
                search_part, purpose_part = rest.split("Why You're Searching", 1)
                result["search_term"] = search_part.strip().split("\n")[0].strip()
                result["search_purpose"] = purpose_part.strip().split("\n\n")[0].strip()
            else:
                result["search_term"] = rest.strip().split("\n")[0].strip()

    # Key insights are the titled blocks after the search info
    # They appear as short title + longer description
    paragraphs = re.split(r'\n\n+', text)
    for para in paragraphs:
        para = para.strip()
        if not para:
            continue
        # Skip the search metadata
        if any(x in para for x in ["What You Searched", "Why You're Searching"]):
            continue
        if len(para) > 30:
            result["key_insights"].append(para)

    if result["key_insights"]:
        result["overview"] = result["key_insights"][0]

    return result


def deinterleave_columns(lines: list[str], n_cols: int) -> list[list[str]]:
    """De-interleave text lines from a multi-column PDF layout.

    First columns may have more lines than later columns due to vertical
    offset in the PDF grid. Pattern: partial first group (remainder columns),
    then full groups of n_cols.
    """
    n_lines = len(lines)
    remainder = n_lines % n_cols
    columns: list[list[str]] = [[] for _ in range(n_cols)]
    idx = 0
    # Partial first group: first `remainder` columns get one extra line
    for col in range(remainder):
        if idx < n_lines:
            columns[col].append(lines[idx])
            idx += 1
    # Full groups
    while idx < n_lines:
        for col in range(n_cols):
            if idx < n_lines:
                columns[col].append(lines[idx])
                idx += 1
    return columns


def parse_top_posts_pair(text: str) -> tuple[dict, dict]:
    """Parse a 2-column Most/Least section into (most_post, least_post).

    Handles variable pdftotext layouts where Most (left) and Least (right)
    column content interleaves. Strategy: first occurrence of each metric
    label → Most, second → Least.
    """
    most = {"caption": "", "engagement_rate": 0, "likes": 0, "comments": 0, "link": ""}
    least = {"caption": "", "engagement_rate": 0, "likes": 0, "comments": 0, "link": ""}

    lines = text.split("\n")

    METRIC_LABELS = {"Caption", "Engagement Rate", "Likes Count", "Comment Count", "Link"}
    HEADER_LABELS = {"Most Liked", "Least Liked", "Most Comments", "Least Comments",
                     "Most Engaged", "Least Engaged"}
    FOOTNOTE_STARTS = ("All performance", "Engagement rate calculated", "Note statistics")

    # Collect labeled sections: list of (label, [value_lines])
    label_groups: list[tuple[str, list[str]]] = []
    i = 0
    while i < len(lines):
        stripped = lines[i].strip()

        if not stripped or stripped in HEADER_LABELS:
            i += 1
            continue

        if any(stripped.startswith(f) for f in FOOTNOTE_STARTS):
            i += 1
            continue

        if stripped in METRIC_LABELS:
            label = stripped
            value_lines = []
            j = i + 1
            while j < len(lines):
                val = lines[j].strip()
                if not val:
                    j += 1
                    continue
                if val in METRIC_LABELS:
                    break
                # Skip header labels (Most/Least markers) — they can appear
                # between a metric label and its value in 2-column layouts
                if val in HEADER_LABELS:
                    j += 1
                    continue
                if any(val.startswith(f) for f in FOOTNOTE_STARTS):
                    break
                value_lines.append(val)
                # For non-Caption metrics, stop after first value line
                if label != "Caption":
                    j += 1
                    break
                j += 1
            label_groups.append((label, value_lines))
            i = j
        else:
            i += 1

    # Map: first occurrence of each label → Most, second → Least
    seen_counts: dict[str, int] = {}
    for label, value_lines in label_groups:
        count = seen_counts.get(label, 0)
        target = most if count == 0 else least

        if label == "Caption":
            target["caption"] = " ".join(value_lines)
        elif label == "Engagement Rate" and value_lines:
            target["engagement_rate"] = parse_number(value_lines[0])
        elif label == "Likes Count" and value_lines:
            target["likes"] = parse_number(value_lines[0])
        elif label == "Comment Count" and value_lines:
            target["comments"] = parse_number(value_lines[0])
        elif label == "Link" and value_lines:
            target["link"] = value_lines[0]

        seen_counts[label] = count + 1

    return most, least


def parse_statistics(text: str) -> dict:
    """Parse summary statistics section."""
    stats = {
        "all_posts": {
            "min_views": 0, "max_views": 0, "median_views": 0, "avg_views": 0,
            "min_likes": 0, "max_likes": 0, "median_likes": 0, "avg_likes": 0,
            "min_comments": 0, "max_comments": 0, "median_comments": 0, "avg_comments": 0,
        }
    }

    lines = text.split("\n")
    # Look for labeled number patterns
    for i, line in enumerate(lines):
        stripped = line.strip()
        label_map = {
            "Min Views": "min_views", "Max Views": "max_views",
            "Median Views": "median_views", "Avg Views": "avg_views",
            "Min Likes": "min_likes", "Max Likes": "max_likes",
            "Median Likes": "median_likes", "Avg Likes": "avg_likes",
            "Min Comments": "min_comments", "Max Comments": "max_comments",
            "Median Comments": "median_comments", "Avg Comments": "avg_comments",
        }
        for label, key in label_map.items():
            if stripped == label:
                for j in range(i + 1, min(i + 3, len(lines))):
                    m = re.search(r'[\d,]+\.?\d*', lines[j])
                    if m:
                        stats["all_posts"][key] = parse_number(m.group())
                        break

    return stats


# ── Main Parsers per Report Type ──────────────────────────────────────

def parse_tiktok_hashtag(sections: dict, identifier: str, date_str: str) -> dict:
    """Build JSON for a TikTok hashtag analysis report."""
    report = {
        "report_type": "tiktok_hashtag",
        "hashtag": identifier,
        "report_date": date_str,
    }

    if "Executive Summary" in sections:
        report["executive_summary"] = parse_executive_summary(sections["Executive Summary"])

    if "Audience Profile" in sections:
        report["audience_profile"] = parse_nopd(sections["Audience Profile"])

    if "Hashtag Analysis" in sections:
        report["hashtag_analysis"] = parse_hashtag_analysis(sections["Hashtag Analysis"])

    if "Interesting Conversations" in sections:
        convs = parse_interesting_conversations(sections["Interesting Conversations"])
        report["interesting_conversations"] = convs

    if "Conversation Map" in sections:
        report["conversation_map"] = parse_conversation_map(sections["Conversation Map"])

    if "Content Trends" in sections:
        report["content_trends"] = parse_content_trends(sections["Content Trends"])

    if "Brand Mentions" in sections:
        report["brand_mentions"] = parse_brand_mentions(sections["Brand Mentions"])

    if "In-Market Campaigns" in sections:
        report["in_market_campaigns"] = parse_in_market_campaigns(
            sections["In-Market Campaigns"])

    if "How to Win With This Audience" in sections:
        report["how_to_win"] = parse_how_to_win(
            sections["How to Win With This Audience"])

    if "Creator Archetypes" in sections:
        report["creator_archetypes"] = parse_creator_archetypes(
            sections["Creator Archetypes"])

    return report


def parse_tiktok_profile(sections: dict, identifier: str, date_str: str) -> dict:
    """Build JSON for a TikTok profile analysis report."""
    report = {
        "report_type": "tiktok_profile",
        "username": identifier,
        "report_date": date_str,
    }

    if "Audience Profile" in sections:
        report["audience_profile"] = parse_nopd(sections["Audience Profile"])

    if "Snapshot" in sections:
        report["snapshot"] = parse_snapshot(sections["Snapshot"])

    if "Creator Summary" in sections:
        report["creator_summary"] = parse_creator_summary(sections["Creator Summary"])

    if "Sponsorship Analysis" in sections:
        report["sponsorships"] = parse_sponsorships(sections["Sponsorship Analysis"])

    if "Future Sponsorship Suggestions" in sections:
        report["future_sponsorship_suggestions"] = parse_future_sponsorships(
            sections["Future Sponsorship Suggestions"])

    if "Engagement Analysis" in sections:
        report["engagement_analysis"] = {"summary": sections["Engagement Analysis"].strip()}

    if "Summary Statistics" in sections:
        report["statistics"] = parse_statistics(sections["Summary Statistics"])
    elif "Summary Statistics - All Posts" in sections:
        report["statistics"] = parse_statistics(sections["Summary Statistics - All Posts"])

    if "Posting Analysis" in sections:
        report["posting_analysis"] = {"summary": sections["Posting Analysis"].strip()}

    if "How to Win With This Audience" in sections:
        report["how_to_win"] = parse_how_to_win(
            sections["How to Win With This Audience"])

    # Top posts — extract both Most and Least from each section
    top_posts = {}
    if "Most / Least Liked" in sections:
        m, l = parse_top_posts_pair(sections["Most / Least Liked"])
        top_posts["most_liked"] = m
        top_posts["least_liked"] = l
    if "Most / Least Comments" in sections:
        m, l = parse_top_posts_pair(sections["Most / Least Comments"])
        top_posts["most_comments"] = m
        top_posts["least_comments"] = l
    if "Most / Least Engaged" in sections:
        m, l = parse_top_posts_pair(sections["Most / Least Engaged"])
        top_posts["most_engaged"] = m
        top_posts["least_engaged"] = l
    if top_posts:
        report["top_posts"] = top_posts

    return report


def parse_instagram_profile(sections: dict, identifier: str, date_str: str) -> dict:
    """Build JSON for an Instagram profile analysis report."""
    report = {
        "report_type": "instagram_profile",
        "username": identifier,
        "report_date": date_str,
    }

    if "Snapshot" in sections:
        report["snapshot"] = parse_snapshot(sections["Snapshot"])

    if "Creator Summary" in sections:
        report["creator_summary"] = parse_creator_summary(sections["Creator Summary"])

    if "Audience Profile" in sections:
        report["audience_profile"] = parse_nopd(sections["Audience Profile"])

    if "Sponsorship Analysis" in sections:
        report["sponsorships"] = parse_sponsorships(sections["Sponsorship Analysis"])

    if "Future Sponsorship Suggestions" in sections:
        report["future_sponsorship_suggestions"] = parse_future_sponsorships(
            sections["Future Sponsorship Suggestions"])

    if "Summary Statistics - All Posts" in sections:
        report["statistics"] = parse_statistics(sections["Summary Statistics - All Posts"])
    elif "Summary Statistics" in sections:
        report["statistics"] = parse_statistics(sections["Summary Statistics"])

    if "Engagement Analysis" in sections:
        report["engagement_analysis"] = {"summary": sections["Engagement Analysis"].strip()}

    if "Posting Analysis" in sections:
        report["posting_analysis"] = {"summary": sections["Posting Analysis"].strip()}

    if "How to Win With This Audience" in sections:
        report["how_to_win"] = parse_how_to_win(
            sections["How to Win With This Audience"])

    # Top posts — extract both Most and Least from each section
    top_posts = {}
    if "Most / Least Liked" in sections:
        m, l = parse_top_posts_pair(sections["Most / Least Liked"])
        top_posts["most_liked"] = m
        top_posts["least_liked"] = l
    if "Most / Least Comments" in sections:
        m, l = parse_top_posts_pair(sections["Most / Least Comments"])
        top_posts["most_comments"] = m
        top_posts["least_comments"] = l
    if "Most / Least Engaged" in sections:
        m, l = parse_top_posts_pair(sections["Most / Least Engaged"])
        top_posts["most_engaged"] = m
        top_posts["least_engaged"] = l
    if top_posts:
        report["top_posts"] = top_posts

    return report


def parse_instagram_hashtag(sections: dict, identifier: str, date_str: str) -> dict:
    """Build JSON for an Instagram hashtag analysis report."""
    # Very similar to TikTok hashtag but without conversation map and in-market campaigns
    report = {
        "report_type": "instagram_hashtag",
        "hashtag": identifier,
        "report_date": date_str,
    }

    if "Executive Summary" in sections:
        report["executive_summary"] = parse_executive_summary(sections["Executive Summary"])

    if "Audience Profile" in sections:
        report["audience_profile"] = parse_nopd(sections["Audience Profile"])

    if "Hashtag Analysis" in sections:
        report["hashtag_analysis"] = parse_hashtag_analysis(sections["Hashtag Analysis"])

    if "Conversation Map" in sections:
        report["conversation_map"] = parse_conversation_map(sections["Conversation Map"])

    if "Content Trends" in sections:
        report["content_trends"] = parse_content_trends(sections["Content Trends"])

    if "Brand Mentions" in sections:
        report["brand_mentions"] = parse_brand_mentions(sections["Brand Mentions"])

    if "Creator Archetypes" in sections:
        report["creator_archetypes"] = parse_creator_archetypes(
            sections["Creator Archetypes"])

    if "How to Win With This Audience" in sections:
        report["how_to_win"] = parse_how_to_win(
            sections["How to Win With This Audience"])

    return report


def parse_tiktok_keywords(sections: dict, identifier: str, date_str: str) -> dict:
    """Build JSON for a TikTok keyword analysis report."""
    report = {
        "report_type": "tiktok_keywords",
        "keyword": identifier,
        "report_date": date_str,
    }

    if "Executive Summary" in sections:
        report["executive_summary"] = parse_executive_summary(sections["Executive Summary"])

    if "Audience Profile" in sections:
        report["audience_profile"] = parse_nopd(sections["Audience Profile"])

    if "Content Trends" in sections:
        report["content_trends"] = parse_content_trends(sections["Content Trends"])

    if "Brand Mentions" in sections:
        report["brand_mentions"] = parse_brand_mentions(sections["Brand Mentions"])

    if "Creator Archetypes" in sections:
        report["creator_archetypes"] = parse_creator_archetypes(
            sections["Creator Archetypes"])

    if "How to Win With This Audience" in sections:
        report["how_to_win"] = parse_how_to_win(
            sections["How to Win With This Audience"])

    return report


def parse_instagram_keywords(sections: dict, identifier: str, date_str: str) -> dict:
    """Build JSON for an Instagram keyword search analysis report.

    Same structure as TikTok keywords but adds Interesting Conversations,
    In-Market Campaigns, and Verbatim sections when present.
    """
    report = {
        "report_type": "instagram_keywords",
        "keyword": identifier,
        "report_date": date_str,
    }

    if "Executive Summary" in sections:
        report["executive_summary"] = parse_executive_summary(sections["Executive Summary"])

    if "Audience Profile" in sections:
        report["audience_profile"] = parse_nopd(sections["Audience Profile"])

    if "Content Trends" in sections:
        report["content_trends"] = parse_content_trends(sections["Content Trends"])

    if "Brand Mentions" in sections:
        report["brand_mentions"] = parse_brand_mentions(sections["Brand Mentions"])

    if "Creator Archetypes" in sections:
        report["creator_archetypes"] = parse_creator_archetypes(
            sections["Creator Archetypes"])

    if "How to Win With This Audience" in sections:
        report["how_to_win"] = parse_how_to_win(
            sections["How to Win With This Audience"])

    if "Interesting Conversations" in sections:
        report["interesting_conversations"] = parse_interesting_conversations(
            sections["Interesting Conversations"])

    if "In-Market Campaigns" in sections:
        report["in_market_campaigns"] = parse_in_market_campaigns(
            sections["In-Market Campaigns"])

    if "Verbatim" in sections:
        # Verbatim is typically raw quotes — store as list of lines
        raw = sections["Verbatim"].strip()
        report["verbatim"] = [ln.strip() for ln in raw.split("\n") if ln.strip()]

    return report


def parse_google_news(sections: dict, identifier: str, date_str: str) -> dict:
    """Build JSON for a Google News analysis report."""
    report = {
        "report_type": "google_news",
        "search_query": identifier,
        "report_date": date_str,
    }

    # ── Executive Summary ──────────────────────────────────────────────
    if "Executive Summary" in sections:
        report["executive_summary"] = parse_executive_summary(sections["Executive Summary"])

    # ── Audience Profile (NOPD) ────────────────────────────────────────
    if "Audience Profile" in sections:
        report["audience_profile"] = parse_nopd(sections["Audience Profile"])

    # ── News Analysis (constructed from multiple sections) ─────────────
    news_analysis = {
        "summary": "",
        "volume_trend": "",
        "sentiment_breakdown": {
            "positive_pct": 0,
            "neutral_pct": 0,
            "negative_pct": 0,
        },
        "key_topics": [],
        "key_findings": [],
        "opportunities": [],
        "risks": [],
    }

    # News Topics → key_topics and key_findings
    if "News Topics" in sections:
        topics_data = parse_news_topics(sections["News Topics"])
        news_analysis["key_topics"] = topics_data["topics"]
        news_analysis["key_findings"] = topics_data["findings"]

    # SWOT Analysis → opportunities and risks + standalone swot_analysis
    if "SWOT Analysis" in sections:
        swot = parse_swot_overview(sections["SWOT Analysis"])
        news_analysis["opportunities"] = swot.get("opportunities", [])
        news_analysis["risks"] = swot.get("threats", [])
        report["swot_analysis"] = swot

    # News Trends → summary for news analysis
    if "News Trends" in sections:
        paragraphs = [p.strip() for p in sections["News Trends"].split("\n\n")
                      if p.strip() and len(p.strip()) > 50]
        if paragraphs:
            news_analysis["summary"] = " ".join(paragraphs[0].split())

    # Fallback: raw News Analysis section text
    if not news_analysis["summary"] and "News Analysis" in sections:
        raw = sections["News Analysis"].strip()
        # Skip boilerplate text about saving hours
        if "saved yourself" not in raw.lower() and len(raw) > 30:
            news_analysis["summary"] = raw

    report["news_analysis"] = news_analysis

    # ── Trending Narratives from News Trends ───────────────────────────
    if "News Trends" in sections:
        report["trending_narratives"] = parse_news_trends(sections["News Trends"])
    elif "Trending Narratives" in sections:
        report["trending_narratives"] = parse_content_trends(sections["Trending Narratives"])

    # ── Brand Mentions ─────────────────────────────────────────────────
    if "Brand Mentions" in sections:
        report["brand_mentions"] = parse_brand_mentions(sections["Brand Mentions"])

    # ── In-Market Campaigns ────────────────────────────────────────────
    if "In-Market Campaigns" in sections:
        report["in_market_campaigns"] = parse_in_market_campaigns(
            sections["In-Market Campaigns"])

    # ── Strategic Implications (from Consideration Spaces + Actions) ───
    strat = {"summary": "", "action_items": []}
    if "Consideration Spaces" in sections:
        # Extract overview paragraph (first long paragraph)
        paras = [p.strip() for p in re.split(r'\n\n+', sections["Consideration Spaces"])
                 if p.strip()]
        long_paras = [p for p in paras if len(p) > 100]
        short_paras = [" ".join(p.split()) for p in paras if 20 < len(p) <= 100]
        if long_paras:
            strat["summary"] = " ".join(long_paras[0].split())
        # Include the consideration area titles as part of the summary
        if short_paras:
            strat["summary"] += "\n\nKey strategic areas: " + " | ".join(short_paras)

    if "Potential Actions" in sections:
        strat["action_items"] = parse_potential_actions(sections["Potential Actions"])

    # Fallback: use Strategic Implications section if present
    if not strat["summary"] and "Strategic Implications" in sections:
        strat["summary"] = sections["Strategic Implications"].strip()

    report["strategic_implications"] = strat

    # ── Quotes ─────────────────────────────────────────────────────────
    if "Quotes" in sections:
        quotes_text = sections["Quotes"]
        # Statistics may be embedded after Quotes (with form-feed or newline before it)
        stats_match = re.search(r'[\n\f]Statistics\n', quotes_text)
        if stats_match:
            report["quotes"] = parse_news_quotes(quotes_text[:stats_match.start()])
            report["key_statistics"] = parse_news_statistics(
                quotes_text[stats_match.end():])
        else:
            report["quotes"] = parse_news_quotes(quotes_text)

    # ── Top Stories (if present as a section) ──────────────────────────
    if "Top Stories" in sections:
        stories = []
        for para in re.split(r'\n\n+', sections["Top Stories"]):
            para = para.strip()
            if para and len(para) > 30:
                stories.append({
                    "headline": " ".join(para.split()),
                    "source": "",
                    "date": "",
                    "summary": "",
                    "sentiment": "",
                    "url": "",
                })
        report["top_stories"] = stories

    # ── Competitor Coverage (if present as a section) ──────────────────
    if "Competitor Coverage" in sections:
        report["competitor_coverage"] = parse_brand_mentions(
            sections["Competitor Coverage"])

    return report


# ── Dispatch ──────────────────────────────────────────────────────────

PARSERS = {
    "tiktok_hashtags": parse_tiktok_hashtag,
    "tiktok_profiles": parse_tiktok_profile,
    "instagram_profiles": parse_instagram_profile,
    "instagram_hashtags": parse_instagram_hashtag,
    "instagram_keywords": parse_instagram_keywords,
    "tiktok_keywords": parse_tiktok_keywords,
    "google_news": parse_google_news,
}


def parse_pdf(pdf_path: str) -> tuple[str, str, dict]:
    """Parse a single PDF or PPTX file into structured JSON.

    Returns (report_type_dir, identifier, report_dict).
    """
    if pdf_path.lower().endswith(".pptx"):
        text = extract_text_from_pptx(pdf_path)
    else:
        text = extract_text_from_pdf(pdf_path)
    report_type_dir, identifier, date_str = detect_report_type(text)
    sections = split_into_sections(text)

    parser_func = PARSERS.get(report_type_dir)
    if not parser_func:
        raise ValueError(f"No parser for report type: {report_type_dir}")

    report = parser_func(sections, identifier, date_str)
    return report_type_dir, identifier, report


def parse_and_save_pdf(pdf_path: str) -> str:
    """Parse a PDF and save the JSON output to the correct directory.

    Returns the output JSON file path.
    """
    report_type_dir, identifier, report = parse_pdf(pdf_path)

    # Create safe filename from identifier
    safe_name = re.sub(r'[^\w\-]', '_', identifier.lower()).strip('_')
    output_dir = os.path.join(AUTOSTRAT_DIR, report_type_dir)
    os.makedirs(output_dir, exist_ok=True)
    output_path = os.path.join(output_dir, f"{safe_name}.json")

    with open(output_path, "w", encoding="utf-8") as f:
        json.dump(report, f, indent=2, ensure_ascii=False)

    return output_path


def parse_all_pdfs(pdf_dir: str = None) -> list[dict]:
    """Parse all PDFs in the given directory.

    Returns list of {pdf: str, output: str, report_type: str, identifier: str, error: str|None}.
    """
    if pdf_dir is None:
        pdf_dir = PDF_DIR

    if not os.path.isdir(pdf_dir):
        return []

    results = []
    for filename in sorted(os.listdir(pdf_dir)):
        if not (filename.lower().endswith(".pdf") or filename.lower().endswith(".pptx")):
            continue

        pdf_path = os.path.join(pdf_dir, filename)
        try:
            report_type_dir, identifier, report = parse_pdf(pdf_path)
            safe_name = re.sub(r'[^\w\-]', '_', identifier.lower()).strip('_')
            output_dir = os.path.join(AUTOSTRAT_DIR, report_type_dir)
            os.makedirs(output_dir, exist_ok=True)
            output_path = os.path.join(output_dir, f"{safe_name}.json")

            with open(output_path, "w", encoding="utf-8") as f:
                json.dump(report, f, indent=2, ensure_ascii=False)

            results.append({
                "pdf": filename,
                "output": output_path,
                "report_type": report_type_dir,
                "identifier": identifier,
                "error": None,
            })
        except Exception as e:
            results.append({
                "pdf": filename,
                "output": None,
                "report_type": None,
                "identifier": None,
                "error": str(e),
            })

    return results
